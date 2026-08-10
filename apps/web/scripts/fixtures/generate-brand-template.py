#!/usr/bin/env python3
"""Generate the fictional "Meridian Advisory" brand template used by the
training-video capture pipeline (capture-deck-frames.cjs).

The brand is invented for this repository so the demo material stays free of
any real company's identity. Rebuild with:

    services/api/.venv/bin/python apps/web/scripts/fixtures/generate-brand-template.py

Produces brand-template.pptx beside this script: four layouts carrying
full-bleed designs (flattened by the server parser), a master logo picture,
patched theme accent colors, and six text slides mapped across the layouts.
"""

from __future__ import annotations

import io
from copy import deepcopy
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches

OUT = Path(__file__).parent / "brand-template.pptx"

INDIGO = (31, 42, 86)
SLATE = (61, 78, 138)
AMBER = (232, 163, 61)
PAPER = (245, 246, 250)

W, H = 1280, 720


def _font(size: int) -> ImageFont.FreeTypeFont:
    try:
        return ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size)
    except OSError:
        return ImageFont.load_default()


def _wordmark(draw: ImageDraw.ImageDraw, x: int, y: int, size: int, color) -> None:
    draw.ellipse([x, y + size * 0.15, x + size * 0.7, y + size * 0.85], outline=color, width=max(2, size // 12))
    draw.text((x + size, y), "MERIDIAN", font=_font(size), fill=color)


def _png(image: Image.Image) -> bytes:
    buf = io.BytesIO()
    image.save(buf, format="PNG")
    return buf.getvalue()


def design_title() -> bytes:
    img = Image.new("RGB", (W, H), INDIGO)
    d = ImageDraw.Draw(img)
    d.polygon([(0, H), (W * 0.42, H), (W * 0.30, H * 0.55), (0, H * 0.72)], fill=SLATE)
    d.rectangle([0, H * 0.86, W, H * 0.885], fill=AMBER)
    _wordmark(d, W - 330, H - 78, 34, PAPER)
    return _png(img)


def design_content() -> bytes:
    img = Image.new("RGB", (W, H), PAPER)
    d = ImageDraw.Draw(img)
    d.rectangle([0, 0, W, 88], fill=INDIGO)
    d.rectangle([0, 88, W, 96], fill=AMBER)
    _wordmark(d, W - 260, H - 56, 26, SLATE)
    return _png(img)


def design_section() -> bytes:
    img = Image.new("RGB", (W, H), SLATE)
    d = ImageDraw.Draw(img)
    d.rectangle([0, 0, W * 0.38, H], fill=INDIGO)
    d.ellipse([W * 0.30, H * 0.32, W * 0.46, H * 0.68], outline=AMBER, width=10)
    _wordmark(d, W - 330, H - 78, 34, PAPER)
    return _png(img)


def design_closing() -> bytes:
    img = Image.new("RGB", (W, H), INDIGO)
    d = ImageDraw.Draw(img)
    d.rectangle([0, H * 0.78, W, H], fill=SLATE)
    d.ellipse([W * 0.44, H * 0.20, W * 0.56, H * 0.42], outline=AMBER, width=8)
    _wordmark(d, W // 2 - 165, int(H * 0.52), 34, PAPER)
    return _png(img)


def logo_png() -> bytes:
    img = Image.new("RGBA", (420, 96), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    _wordmark(d, 4, 14, 52, INDIGO)
    return _png(img)


def _bleed_layout(presentation: Presentation, index: int, art: bytes) -> None:
    """Draw full-bleed artwork onto a layout part (python-pptx has no
    high-level API for layout pictures: draw on a throwaway slide, then
    re-relate the picture element to the layout)."""
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # blank
    picture = slide.shapes.add_picture(
        io.BytesIO(art), 0, 0, width=presentation.slide_width, height=presentation.slide_height
    )
    blip = picture._element.find(f".//{qn('a:blip')}")
    r_id = blip.get(qn("r:embed"))
    image_part = slide.part.related_part(r_id)

    layout = presentation.slide_layouts[index]
    layout_rid = layout.part.relate_to(image_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")
    layout_pic = deepcopy(picture._element)
    layout_pic.find(f".//{qn('a:blip')}").set(qn("r:embed"), layout_rid)
    # Put the art behind the layout's placeholders.
    layout.shapes._spTree.insert(2, layout_pic)

    picture._element.getparent().remove(picture._element)
    _drop_slide(presentation, slide)


def _drop_slide(presentation: Presentation, slide) -> None:
    xml_slides = presentation.slides._sldIdLst
    for sld_id in list(xml_slides):
        if presentation.part.related_part(sld_id.get(qn("r:id"))) is slide.part:
            xml_slides.remove(sld_id)


def _master_logo(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    picture = slide.shapes.add_picture(
        io.BytesIO(logo_png()), Emu(0), Emu(0), width=Inches(1.9), height=Inches(0.44)
    )
    picture.left = presentation.slide_width - picture.width - Inches(0.35)
    picture.top = presentation.slide_height - picture.height - Inches(0.25)
    blip = picture._element.find(f".//{qn('a:blip')}")
    image_part = slide.part.related_part(blip.get(qn("r:embed")))

    master = presentation.slide_masters[0]
    master_rid = master.part.relate_to(image_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")
    master_pic = deepcopy(picture._element)
    master_pic.find(f".//{qn('a:blip')}").set(qn("r:embed"), master_rid)
    master.shapes._spTree.append(master_pic)

    picture._element.getparent().remove(picture._element)
    _drop_slide(presentation, slide)


def _patch_theme_colors(presentation: Presentation) -> None:
    theme = presentation.slide_masters[0].element.getroottree().getroot()  # placeholder; real theme below
    part = presentation.slide_masters[0].part
    for rel in part.rels.values():
        if rel.reltype.endswith("/theme"):
            theme_part = rel.target_part
            break
    else:
        return
    root = theme_part._element if hasattr(theme_part, "_element") else None
    import lxml.etree as etree

    root = etree.fromstring(theme_part.blob)
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for name, rgb in (
        ("dk2", "1F2A56"),
        ("lt2", "F5F6FA"),
        ("accent1", "1F2A56"),
        ("accent2", "E8A33D"),
        ("accent3", "3D4E8A"),
        ("accent4", "8A96C4"),
        ("accent5", "C9A227"),
        ("accent6", "0F1530"),
    ):
        el = root.find(f".//a:clrScheme/a:{name}/a:srgbClr", ns)
        if el is not None:
            el.set("val", rgb)
    theme_part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def main() -> None:
    presentation = Presentation()

    designs = {0: design_title(), 1: design_content(), 2: design_section(), 5: design_closing()}
    for index, art in designs.items():
        _bleed_layout(presentation, index, art)
    _master_logo(presentation)
    _patch_theme_colors(presentation)

    def add(layout_index: int, title: str, body_lines: list[str] | None = None) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
        if body_lines:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    frame = shape.text_frame
                    frame.text = body_lines[0]
                    for line in body_lines[1:]:
                        frame.add_paragraph().text = line
                    break

    add(0, "Meridian Advisory", ["Client workshop template"])
    add(1, "Agenda", ["Where we are today", "What the data shows", "Options on the table", "Recommended path"])
    add(2, "Part One - Discovery")
    add(1, "Our Approach", ["Interview the working team", "Map the current process", "Quantify the gaps", "Design the target state"])
    add(1, "Timeline", ["Weeks 1-2: discovery", "Weeks 3-4: analysis", "Weeks 5-6: recommendations", "Week 7: readout"])
    add(5, "Thank You", ["meridian.example.com"])

    presentation.save(OUT)
    print(f"wrote {OUT} ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
