from __future__ import annotations

import base64
import io

from fastapi.testclient import TestClient
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches

from app.core.deck_template_parse import parse_deck_template
from app.main import app

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
AUTH = {"x-aperture-user": "user-admin"}

# A 1x1 PNG for embedding as a small (logo-sized) picture.
_TINY_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
    "0000000d49444154789c626001000000ffff030000060005579bcaa30000000049454e44ae426082"
)


def _template_bytes(*, with_logo: bool = False) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Brand Deck"
    slide.placeholders[1].text = "Subtitle scaffolding"
    body_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    body_slide.shapes.title.text = "Agenda"
    body_slide.placeholders[1].text_frame.text = "First point"
    if with_logo:
        # python-pptx has no high-level API for master/layout pictures (real
        # templates get them from PowerPoint), so author the XML directly:
        # draw the picture on a slide, then re-relate it to a layout part.
        from copy import deepcopy

        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.oxml.ns import qn

        picture = slide.shapes.add_picture(
            io.BytesIO(_TINY_PNG), Inches(0.2), Inches(0.2), Inches(0.5), Inches(0.5)
        )
        blip = picture._element.find(f".//{qn('a:blip')}")
        old_rid = blip.get(qn("r:embed"))
        image_part = slide.part.related_part(old_rid)
        layout = presentation.slide_layouts[0]
        new_rid = layout.part.relate_to(image_part, RT.IMAGE)
        layout_pic = deepcopy(picture._element)
        layout_pic.find(f".//{qn('a:blip')}").set(qn("r:embed"), new_rid)
        layout.shapes._spTree.append(layout_pic)
        picture._element.getparent().remove(picture._element)
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_parse_extracts_theme_fonts_and_slide_text() -> None:
    result = parse_deck_template("brand.pptx", _template_bytes())

    assert result.slide_count == 2
    # Stock Office theme colors round-trip as hex.
    assert result.theme.colors["lt1"] == "#ffffff"
    assert result.theme.colors["accent1"].startswith("#")
    assert result.theme.major_font
    titles = [slide.title for slide in result.slides]
    assert "Brand Deck" in titles
    assert any("First point" in " ".join(slide.blocks) for slide in result.slides)


def test_parse_finds_master_logo_candidates() -> None:
    result = parse_deck_template("brand.pptx", _template_bytes(with_logo=True))

    assert result.logo_candidates
    assert result.logo_candidates[0].data_url.startswith("data:image/png;base64,")
    assert result.logo_candidates[0].source.startswith(("master", "layout:"))


def test_endpoint_parses_uploads_and_rejects_unsafe_files() -> None:
    client = TestClient(app)

    ok = client.post(
        "/api/drafts/deck-template/parse",
        files={"file": ("brand.pptx", _template_bytes(), PPTX_MIME)},
        headers=AUTH,
    )
    assert ok.status_code == 200
    payload = ok.json()
    assert payload["filename"] == "brand.pptx"
    assert payload["theme"]["colors"]

    macro = client.post(
        "/api/drafts/deck-template/parse",
        files={"file": ("brand.pptm", _template_bytes(), PPTX_MIME)},
        headers=AUTH,
    )
    assert macro.status_code == 415
    assert "Macro-enabled" in macro.json()["detail"]

    wrong_type = client.post(
        "/api/drafts/deck-template/parse",
        files={"file": ("notes.txt", b"hello", "text/plain")},
        headers=AUTH,
    )
    assert wrong_type.status_code == 415

    not_zip = client.post(
        "/api/drafts/deck-template/parse",
        files={"file": ("fake.pptx", b"not a zip at all", PPTX_MIME)},
        headers=AUTH,
    )
    assert not_zip.status_code == 422

    corrupt = client.post(
        "/api/drafts/deck-template/parse",
        files={"file": ("corrupt.pptx", b"PK\x03\x04garbage-that-is-not-a-package", PPTX_MIME)},
        headers=AUTH,
    )
    assert corrupt.status_code == 422


def test_endpoint_requires_authentication() -> None:
    client = TestClient(app)
    response = client.post(
        "/api/drafts/deck-template/parse",
        files={"file": ("brand.pptx", _template_bytes(), PPTX_MIME)},
    )
    assert response.status_code in (401, 403)


def _multi_design_template_bytes() -> bytes:
    """Two slides on two different layouts, each layout carrying its own
    full-bleed coloured picture — the shape a real brand template has."""

    from copy import deepcopy

    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import qn

    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[0])
    first.shapes.title.text = "Cover"
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    second.shapes.title.text = "Agenda"

    for index, (slide, colour) in enumerate(((first, (200, 30, 30)), (second, (30, 60, 200)))):
        buffer = io.BytesIO()
        PILImage.new("RGB", (160, 90), colour).save(buffer, format="PNG")
        buffer.seek(0)
        picture = slide.shapes.add_picture(
            buffer, 0, 0, presentation.slide_width, presentation.slide_height
        )
        blip = picture._element.find(f".//{qn('a:blip')}")
        image_part = slide.part.related_part(blip.get(qn("r:embed")))
        layout = presentation.slide_layouts[index]
        new_rid = layout.part.relate_to(image_part, RT.IMAGE)
        layout_pic = deepcopy(picture._element)
        layout_pic.find(f".//{qn('a:blip')}").set(qn("r:embed"), new_rid)
        # Layout art sits behind the placeholders.
        layout.shapes._spTree.insert(2, layout_pic)
        picture._element.getparent().remove(picture._element)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_parse_renders_one_design_per_layout_and_maps_slides_to_them() -> None:
    result = parse_deck_template("brand.pptx", _multi_design_template_bytes())

    # Each slide uses a different layout, so each gets its own design.
    assert len(result.designs) == 2
    indexes = [slide.design_index for slide in result.slides]
    assert indexes == [0, 1]
    assert [slide.layout_name for slide in result.slides] == ["Title Slide", "Title and Content"]
    for design in result.designs:
        assert design.data_url.startswith("data:image/jpeg;base64,")
        assert design.width_px == 1280 and design.height_px == 720

    # The rendered designs differ — the whole point of per-slide artwork.
    assert result.designs[0].data_url != result.designs[1].data_url

    # And the flattened art really carries the layout's colour.
    first = PILImage.open(
        io.BytesIO(base64.b64decode(result.designs[0].data_url.split(",", 1)[1]))
    )
    red, green, blue = first.convert("RGB").getpixel((640, 360))
    assert red > 150 and green < 90 and blue < 90


def test_parse_skips_slide_number_placeholders_in_lifted_text() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Agenda"
    slide.placeholders[1].text_frame.text = "Real content"
    # Slide-number furniture must not become slide text.
    from pptx.oxml.ns import qn

    number = slide.shapes.add_textbox(0, 0, 100, 100)
    number.text_frame.text = "12"
    nv_sp_pr = number._element.find(qn("p:nvSpPr"))
    nv_pr = nv_sp_pr.find(qn("p:nvPr"))
    ph = nv_pr.makeelement(qn("p:ph"), {"type": "sldNum", "sz": "quarter", "idx": "12"})
    nv_pr.append(ph)

    buffer = io.BytesIO()
    presentation.save(buffer)
    result = parse_deck_template("furniture.pptx", buffer.getvalue())

    blocks = " ".join(block for slide_text in result.slides for block in slide_text.blocks)
    assert "Real content" in blocks
    assert "12" not in blocks
