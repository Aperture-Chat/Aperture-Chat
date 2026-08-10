"""Brand extraction from uploaded PowerPoint templates.

python-pptx exposes slide structure but has no first-class theme-color API, so
colors and fonts are read straight from the theme part's XML. Every extraction
path degrades to a warning plus usable defaults — a malformed corporate
template must never surface as a 500.
"""

from __future__ import annotations

import base64
import io
import re
from typing import Any

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Emu

from app.models.decks import (
    DeckTemplateImageCandidate,
    DeckTemplateParseResponse,
    DeckTemplateSlideText,
    DeckTemplateTheme,
)

_DRAWINGML_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
_THEME_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
_COLOR_ROLES = ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3", "accent4", "accent5", "accent6")
_ALLOWED_IMAGE_TYPES = {"image/png": "png", "image/jpeg": "jpeg"}
_MAX_IMAGE_BYTES = 512_000
_MAX_BACKGROUND_BYTES = 1_500_000
_MAX_CANDIDATES = 5
_MAX_SLIDES = 60
_MAX_BLOCKS_PER_SLIDE = 12
_MAX_BLOCK_CHARS = 400
_HEX_COLOR = re.compile(r"^[0-9A-Fa-f]{6}$")
# Flattened per-layout design renders. 1280x720 keeps a 16:9 slide sharp on the
# 960x540 stage and in the export while staying small enough that a 30-design
# template still fits the deck's content budget.
_DESIGN_WIDTH = 1280
_DESIGN_HEIGHT = 720
_DESIGN_JPEG_QUALITY = 72
_MAX_DESIGN_BYTES = 900_000
_MAX_DESIGNS = 30
_FURNITURE_PLACEHOLDERS = {
    PP_PLACEHOLDER.SLIDE_NUMBER,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.DATE,
}


def _theme_element(presentation: Any) -> etree._Element | None:
    try:
        master_part = presentation.slide_masters[0].part
        theme_part = master_part.part_related_by(_THEME_RELTYPE)
        return etree.fromstring(theme_part.blob)
    except Exception:
        return None


def _scheme_color(scheme: etree._Element, role: str) -> str | None:
    nodes = scheme.findall(f"a:{role}", _DRAWINGML_NS)
    if not nodes:
        return None
    node = nodes[0]
    srgb = node.find("a:srgbClr", _DRAWINGML_NS)
    if srgb is not None:
        value = srgb.get("val", "")
        return f"#{value.lower()}" if _HEX_COLOR.match(value) else None
    sys_color = node.find("a:sysClr", _DRAWINGML_NS)
    if sys_color is not None:
        value = sys_color.get("lastClr", "")
        return f"#{value.lower()}" if _HEX_COLOR.match(value) else None
    return None


def _extract_theme(presentation: Any, warnings: list[str]) -> DeckTemplateTheme:
    root = _theme_element(presentation)
    if root is None:
        warnings.append("No readable theme part was found; default colors are returned.")
        return DeckTemplateTheme(colors={}, major_font=None, minor_font=None)
    colors: dict[str, str] = {}
    scheme = root.find(".//a:clrScheme", _DRAWINGML_NS)
    if scheme is not None:
        for role in _COLOR_ROLES:
            value = _scheme_color(scheme, role)
            if value:
                colors[role] = value
    else:
        warnings.append("The theme has no color scheme; default colors are returned.")
    major_font: str | None = None
    minor_font: str | None = None
    major = root.find(".//a:fontScheme/a:majorFont/a:latin", _DRAWINGML_NS)
    minor = root.find(".//a:fontScheme/a:minorFont/a:latin", _DRAWINGML_NS)
    if major is not None and (major.get("typeface") or "").strip():
        major_font = major.get("typeface", "").strip()
    if minor is not None and (minor.get("typeface") or "").strip():
        minor_font = minor.get("typeface", "").strip()
    if major_font is None and minor_font is None:
        warnings.append("The theme reports no fonts; the platform default font is used.")
    return DeckTemplateTheme(colors=colors, major_font=major_font, minor_font=minor_font)


def _iter_shapes(shapes: Any):
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) is not None and getattr(shape, "shapes", None):
            yield from _iter_shapes(shape.shapes)


def _picture_candidates(presentation: Any, warnings: list[str]) -> tuple[
    list[DeckTemplateImageCandidate], list[DeckTemplateImageCandidate]
]:
    """Small pictures on masters/layouts are logo candidates; pictures covering
    most of the slide are background candidates."""

    logos: list[DeckTemplateImageCandidate] = []
    backgrounds: list[DeckTemplateImageCandidate] = []
    slide_width = int(presentation.slide_width or Emu(12192000))
    slide_height = int(presentation.slide_height or Emu(6858000))
    slide_area = max(1, slide_width * slide_height)
    seen_digests: set[bytes] = set()

    containers: list[tuple[str, Any]] = []
    for master_index, master in enumerate(presentation.slide_masters, start=1):
        containers.append((f"master{master_index}", master))
        for layout in master.slide_layouts:
            containers.append((f"layout:{layout.name or 'unnamed'}", layout))

    for source, container in containers:
        for shape in _iter_shapes(container.shapes):
            image = getattr(shape, "image", None)
            if image is None:
                continue
            try:
                blob = image.blob
                content_type = image.content_type
                digest = image.sha1.encode() if isinstance(image.sha1, str) else bytes(image.sha1)
            except Exception:
                continue
            if content_type not in _ALLOWED_IMAGE_TYPES or digest in seen_digests:
                continue
            width = int(getattr(shape, "width", 0) or 0)
            height = int(getattr(shape, "height", 0) or 0)
            area_ratio = (width * height) / slide_area if width and height else 0.0
            data_url = (
                f"data:{content_type};base64,{base64.b64encode(blob).decode('ascii')}"
            )
            candidate = DeckTemplateImageCandidate(
                data_url=data_url,
                width_px=max(1, round(Emu(width or slide_width).inches * 96)),
                height_px=max(1, round(Emu(height or slide_height).inches * 96)),
                source=source,
            )
            if area_ratio >= 0.5 and len(blob) <= _MAX_BACKGROUND_BYTES:
                if len(backgrounds) < _MAX_CANDIDATES:
                    backgrounds.append(candidate)
                    seen_digests.add(digest)
            elif area_ratio < 0.25 and len(blob) <= _MAX_IMAGE_BYTES:
                if len(logos) < _MAX_CANDIDATES:
                    logos.append(candidate)
                    seen_digests.add(digest)
    if not logos and not backgrounds:
        warnings.append("No logo or background images were found on the template's masters.")
    return logos, backgrounds


def _emu_box(shape: Any, slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    scale_x = _DESIGN_WIDTH / max(1, slide_width)
    scale_y = _DESIGN_HEIGHT / max(1, slide_height)
    left = int(round((shape.left or 0) * scale_x))
    top = int(round((shape.top or 0) * scale_y))
    width = max(1, int(round((shape.width or 0) * scale_x)))
    height = max(1, int(round((shape.height or 0) * scale_y)))
    return left, top, width, height


def _solid_fill_rgba(shape: Any) -> tuple[int, int, int, int] | None:
    """Explicit sRGB solid fills only. Theme-referenced and gradient fills are
    skipped rather than guessed at — a wrong colour block is worse than none."""

    try:
        if shape.fill.type != MSO_FILL.SOLID:
            return None
        rgb = shape.fill.fore_color.rgb
    except Exception:
        return None
    if rgb is None:
        return None
    try:
        return (rgb[0], rgb[1], rgb[2], 255)
    except Exception:
        return None


def _paste_picture(canvas: Any, shape: Any, slide_width: int, slide_height: int) -> bool:
    try:
        blob = shape.image.blob
        image = Image.open(io.BytesIO(blob))
        image.load()
    except Exception:
        return False
    crop = (
        float(getattr(shape, "crop_left", 0) or 0),
        float(getattr(shape, "crop_top", 0) or 0),
        float(getattr(shape, "crop_right", 0) or 0),
        float(getattr(shape, "crop_bottom", 0) or 0),
    )
    if any(crop):
        width_px, height_px = image.size
        box = (
            int(width_px * crop[0]),
            int(height_px * crop[1]),
            int(width_px * (1 - crop[2])),
            int(height_px * (1 - crop[3])),
        )
        if box[2] > box[0] and box[3] > box[1]:
            image = image.crop(box)
    left, top, width, height = _emu_box(shape, slide_width, slide_height)
    try:
        image = image.convert("RGBA").resize((width, height), Image.LANCZOS)
    except Exception:
        return False
    rotation = float(getattr(shape, "rotation", 0) or 0)
    if rotation:
        image = image.rotate(-rotation, expand=True, resample=Image.BICUBIC)
        left -= (image.width - width) // 2
        top -= (image.height - height) // 2
    canvas.alpha_composite(image, (left, top))
    return True


def _composite_shapes(canvas: Any, shapes: Any, slide_width: int, slide_height: int, depth: int = 0) -> int:
    """Draws a container's art in z-order. Groups recurse; text is deliberately
    left out because the deck editor owns the words."""

    painted = 0
    if depth > 4:
        return painted
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                painted += _composite_shapes(canvas, shape.shapes, slide_width, slide_height, depth + 1)
                continue
            if getattr(shape, "image", None) is not None:
                if _paste_picture(canvas, shape, slide_width, slide_height):
                    painted += 1
                continue
            fill = _solid_fill_rgba(shape)
            if fill is None:
                continue
            left, top, width, height = _emu_box(shape, slide_width, slide_height)
            canvas.alpha_composite(Image.new("RGBA", (width, height), fill), (left, top))
            painted += 1
        except Exception:
            continue
    return painted


def _background_rgb(container: Any, fallback: tuple[int, int, int]) -> tuple[int, int, int]:
    """Reads an explicit <p:bg> solid fill; anything else uses the theme's
    light background."""

    try:
        background = container.element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld"
        ).find("{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
    except Exception:
        return fallback
    if background is None:
        return fallback
    srgb = background.find(".//a:srgbClr", _DRAWINGML_NS)
    if srgb is None:
        return fallback
    value = srgb.get("val", "")
    if not _HEX_COLOR.match(value):
        return fallback
    return (int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _layout_designs(
    presentation: Any, warnings: list[str]
) -> tuple[list[DeckTemplateImageCandidate], dict[int, int]]:
    """Flattens each slide layout's artwork into one background picture.

    A brand template carries its design as pictures and colour blocks on many
    different layouts, so a single extracted image cannot represent the deck.
    Each layout is composited once and reused by every slide that references
    it, which also keeps the payload small. Text is not drawn: it stays
    editable in the deck editor.
    """

    designs: list[DeckTemplateImageCandidate] = []
    by_layout: dict[int, int] = {}
    slide_width = int(presentation.slide_width or 12192000)
    slide_height = int(presentation.slide_height or 6858000)

    used_layouts: list[Any] = []
    for index, slide in enumerate(presentation.slides):
        if index >= _MAX_SLIDES:
            break
        try:
            layout = slide.slide_layout
        except Exception:
            continue
        key = id(layout._element)
        if key not in {id(item._element) for item in used_layouts}:
            used_layouts.append(layout)

    if len(used_layouts) > _MAX_DESIGNS:
        warnings.append(
            f"Only the first {_MAX_DESIGNS} slide designs were rendered; later slides reuse the deck background."
        )
        used_layouts = used_layouts[:_MAX_DESIGNS]

    for layout in used_layouts:
        try:
            master_rgb = _background_rgb(layout.slide_master, (255, 255, 255))
            base_rgb = _background_rgb(layout, master_rgb)
            canvas = Image.new("RGBA", (_DESIGN_WIDTH, _DESIGN_HEIGHT), base_rgb + (255,))
            painted = _composite_shapes(canvas, layout.slide_master.shapes, slide_width, slide_height)
            painted += _composite_shapes(canvas, layout.shapes, slide_width, slide_height)
            if not painted and base_rgb == (255, 255, 255):
                # Nothing to show: let the slide fall back to the theme colour
                # instead of shipping a blank white picture.
                continue
            flattened = canvas.convert("RGB")
            buffer = io.BytesIO()
            flattened.save(buffer, format="JPEG", quality=_DESIGN_JPEG_QUALITY, optimize=True)
            data = buffer.getvalue()
            is_dark = _design_is_dark(flattened)
        except Exception:
            warnings.append(f"The '{getattr(layout, 'name', 'unnamed')}' layout design could not be rendered.")
            continue
        if len(data) > _MAX_DESIGN_BYTES:
            warnings.append(
                f"The '{getattr(layout, 'name', 'unnamed')}' layout design was too large to include."
            )
            continue
        by_layout[id(layout._element)] = len(designs)
        designs.append(
            DeckTemplateImageCandidate(
                data_url=f"data:image/jpeg;base64,{base64.b64encode(data).decode('ascii')}",
                width_px=_DESIGN_WIDTH,
                height_px=_DESIGN_HEIGHT,
                source=f"layout:{getattr(layout, 'name', 'unnamed')}",
                is_dark=is_dark,
            )
        )
    return designs, by_layout


def _design_is_dark(image: Any) -> bool:
    """Mean luminance of the band where slide text sits. Dark artwork needs
    light type, and the client uses this to keep imported slides readable
    instead of printing the brand's dark heading colour onto a photo."""

    try:
        band = image.crop(
            (
                int(_DESIGN_WIDTH * 0.08),
                int(_DESIGN_HEIGHT * 0.18),
                int(_DESIGN_WIDTH * 0.92),
                int(_DESIGN_HEIGHT * 0.82),
            )
        ).resize((32, 18))
        pixels = list(band.getdata())
    except Exception:
        return False
    if not pixels:
        return False
    total = sum(0.2126 * red + 0.7152 * green + 0.0722 * blue for red, green, blue in pixels)
    return (total / len(pixels)) < 128


def _is_furniture_placeholder(shape: Any) -> bool:
    """Slide numbers, footers, and dates are deck furniture, not slide content —
    lifting them turns a template's page numbers into slide titles."""

    try:
        if not shape.is_placeholder:
            return False
        return shape.placeholder_format.type in _FURNITURE_PLACEHOLDERS
    except Exception:
        return False


def _slide_texts(presentation: Any, design_indexes: dict[int, int]) -> list[DeckTemplateSlideText]:
    slides: list[DeckTemplateSlideText] = []
    for index, slide in enumerate(presentation.slides):
        if index >= _MAX_SLIDES:
            break
        title: str | None = None
        blocks: list[str] = []
        try:
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                text = title_shape.text_frame.text.strip()
                title = text[:_MAX_BLOCK_CHARS] or None
        except Exception:
            title = None
        for shape in _iter_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            if title is not None and getattr(slide.shapes, "title", None) is shape:
                continue
            if _is_furniture_placeholder(shape):
                continue
            text = " ".join(shape.text_frame.text.split()).strip()
            if not text or text == title:
                continue
            blocks.append(text[:_MAX_BLOCK_CHARS])
            if len(blocks) >= _MAX_BLOCKS_PER_SLIDE:
                break
        layout_name: str | None = None
        design_index: int | None = None
        try:
            layout = slide.slide_layout
            layout_name = (layout.name or "").strip()[:120] or None
            design_index = design_indexes.get(id(layout._element))
        except Exception:
            layout_name = None
        slides.append(
            DeckTemplateSlideText(
                index=index,
                title=title,
                blocks=blocks,
                layout_name=layout_name,
                design_index=design_index,
            )
        )
    return slides


def parse_deck_template(filename: str, payload: bytes) -> DeckTemplateParseResponse:
    """Parses an uploaded template. Raises ValueError when the file is not a
    readable PowerPoint package; every partial failure becomes a warning."""

    try:
        presentation = Presentation(io.BytesIO(payload))
    except Exception as exc:  # noqa: BLE001 — normalized for the route's 422
        raise ValueError("The file could not be read as a PowerPoint template.") from exc

    warnings: list[str] = []
    theme = _extract_theme(presentation, warnings)
    logos, backgrounds = _picture_candidates(presentation, warnings)
    designs, design_indexes = _layout_designs(presentation, warnings)
    slides = _slide_texts(presentation, design_indexes)
    if slides and not any(slide.design_index is not None for slide in slides):
        warnings.append("No slide designs could be rendered; slides use the extracted brand colors.")
    return DeckTemplateParseResponse(
        filename=filename,
        slide_count=len(slides),
        theme=theme,
        logo_candidates=logos,
        background_candidates=backgrounds,
        designs=designs,
        slides=slides,
        warnings=warnings,
    )
