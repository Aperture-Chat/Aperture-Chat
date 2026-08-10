from __future__ import annotations

import html
import mimetypes
import re
import zipfile
from dataclasses import dataclass
from email import policy
from email.parser import BytesParser
from io import BytesIO
from typing import BinaryIO

from defusedxml import ElementTree
from defusedxml.common import DefusedXmlException

MAX_EXTRACTED_TEXT_CHARS = 10_000_000
DEFAULT_CHUNK_CHARS = 1200
DEFAULT_CHUNK_OVERLAP = 160
DEFAULT_OCR_MAX_PAGES = 250
DEFAULT_OCR_PAGE_TIMEOUT_SECONDS = 45.0
MIN_PDF_PAGE_TEXT_CHARS = 40

_TEXT_EXTENSIONS = {
    ".csv",
    ".json",
    ".log",
    ".md",
    ".rtf",
    ".txt",
    ".xml",
}
_HTML_EXTENSIONS = {".htm", ".html"}
_IMAGE_EXTENSIONS = {".bmp", ".gif", ".jpeg", ".jpg", ".png", ".tif", ".tiff", ".webp"}

_DOCX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_EML_MIME_TYPES = {"message/rfc822", "application/eml"}
_MSG_MIME_TYPES = {"application/vnd.ms-outlook", "application/x-msg"}

BinarySource = bytes | BinaryIO


@dataclass(frozen=True, slots=True)
class ExtractedSegment:
    """One ordered, independently locatable unit of extracted source text.

    ``page_start`` and ``page_end`` are one-based physical PDF/image page
    numbers. ``locator`` carries other honest source positions such as a slide
    or worksheet name. Callers must leave fields empty when the source format
    does not expose a reliable location.
    """

    text: str
    page_start: int | None = None
    page_end: int | None = None
    locator: str | None = None


def extract_text(
    filename: str,
    content: bytes,
    mime_type: str | None = None,
    *,
    max_chars: int = MAX_EXTRACTED_TEXT_CHARS,
    ocr_enabled: bool = True,
    ocr_max_pages: int = DEFAULT_OCR_MAX_PAGES,
    ocr_page_timeout_seconds: float = DEFAULT_OCR_PAGE_TIMEOUT_SECONDS,
) -> str | None:
    return _segments_text(
        extract_segments(
            filename,
            content,
            mime_type,
            max_chars=max_chars,
            ocr_enabled=ocr_enabled,
            ocr_max_pages=ocr_max_pages,
            ocr_page_timeout_seconds=ocr_page_timeout_seconds,
        )
    )


def extract_segments(
    filename: str,
    content: bytes,
    mime_type: str | None = None,
    *,
    max_chars: int = MAX_EXTRACTED_TEXT_CHARS,
    ocr_enabled: bool = True,
    ocr_max_pages: int = DEFAULT_OCR_MAX_PAGES,
    ocr_page_timeout_seconds: float = DEFAULT_OCR_PAGE_TIMEOUT_SECONDS,
) -> list[ExtractedSegment]:
    return _extract_segments_from_source(
        filename,
        content,
        mime_type,
        max_chars=max_chars,
        ocr_enabled=ocr_enabled,
        ocr_max_pages=ocr_max_pages,
        ocr_page_timeout_seconds=ocr_page_timeout_seconds,
    )


def extract_text_from_file(
    filename: str,
    file: BinaryIO,
    mime_type: str | None = None,
    *,
    max_chars: int = MAX_EXTRACTED_TEXT_CHARS,
    ocr_enabled: bool = True,
    ocr_max_pages: int = DEFAULT_OCR_MAX_PAGES,
    ocr_page_timeout_seconds: float = DEFAULT_OCR_PAGE_TIMEOUT_SECONDS,
) -> str | None:
    return _segments_text(
        extract_segments_from_file(
            filename,
            file,
            mime_type,
            max_chars=max_chars,
            ocr_enabled=ocr_enabled,
            ocr_max_pages=ocr_max_pages,
            ocr_page_timeout_seconds=ocr_page_timeout_seconds,
        )
    )


def extract_segments_from_file(
    filename: str,
    file: BinaryIO,
    mime_type: str | None = None,
    *,
    max_chars: int = MAX_EXTRACTED_TEXT_CHARS,
    ocr_enabled: bool = True,
    ocr_max_pages: int = DEFAULT_OCR_MAX_PAGES,
    ocr_page_timeout_seconds: float = DEFAULT_OCR_PAGE_TIMEOUT_SECONDS,
) -> list[ExtractedSegment]:
    return _extract_segments_from_source(
        filename,
        file,
        mime_type,
        max_chars=max_chars,
        ocr_enabled=ocr_enabled,
        ocr_max_pages=ocr_max_pages,
        ocr_page_timeout_seconds=ocr_page_timeout_seconds,
    )


def _extract_segments_from_source(
    filename: str,
    source: BinarySource,
    mime_type: str | None,
    *,
    max_chars: int,
    ocr_enabled: bool,
    ocr_max_pages: int,
    ocr_page_timeout_seconds: float,
) -> list[ExtractedSegment]:
    extension = _extension(filename)
    guessed_mime = (mime_type or mimetypes.guess_type(filename)[0] or "").split(";", 1)[0]
    guessed_mime = guessed_mime.strip().lower()

    segments: list[ExtractedSegment]
    if extension == ".docx" or guessed_mime == _DOCX_MIME_TYPE:
        segments = _extract_docx_segments(source)
    elif extension == ".xlsx" or guessed_mime == _XLSX_MIME_TYPE:
        segments = _extract_xlsx_segments(source)
    elif extension == ".pptx" or guessed_mime == _PPTX_MIME_TYPE:
        segments = _extract_pptx_segments(source)
    elif extension == ".eml" or guessed_mime in _EML_MIME_TYPES:
        segments = _extract_eml_segments(source)
    elif extension == ".msg" or guessed_mime in _MSG_MIME_TYPES:
        segments = _extract_msg_segments(source)
    elif extension in _HTML_EXTENSIONS or guessed_mime == "text/html":
        content = _read_source_bytes(source, max_chars=max_chars)
        text = _extract_html_text(_decode_text(content))
        segments = [ExtractedSegment(text=text)] if text else []
    elif extension == ".pdf" or guessed_mime == "application/pdf":
        segments = _extract_pdf_segments(
            source,
            ocr_enabled=ocr_enabled,
            ocr_max_pages=ocr_max_pages,
            ocr_page_timeout_seconds=ocr_page_timeout_seconds,
        )
    elif extension in _IMAGE_EXTENSIONS or guessed_mime.startswith("image/"):
        if not ocr_enabled:
            segments = []
        else:
            segments = _extract_image_segments(
                source,
                ocr_max_pages=ocr_max_pages,
                ocr_page_timeout_seconds=ocr_page_timeout_seconds,
            )
    elif (
        extension in _TEXT_EXTENSIONS
        or guessed_mime.startswith("text/")
        or guessed_mime in {"application/json", "application/xml"}
    ):
        content = _read_source_bytes(source, max_chars=max_chars)
        text = _decode_text(content)
        if extension == ".rtf":
            text = _extract_rtf_text(text)
        segments = [ExtractedSegment(text=text)] if text else []
    else:
        segments = []
    return _limit_segments(segments, max_chars=max_chars)


def chunk_text(
    text: str,
    *,
    max_chars: int = DEFAULT_CHUNK_CHARS,
    overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> list[str]:
    normalized = _normalize_text(text)
    if not normalized:
        return []
    safe_max = max(200, max_chars)
    safe_overlap = max(0, min(overlap, safe_max // 3))
    paragraphs = [
        paragraph.strip() for paragraph in re.split(r"\n{2,}", normalized) if paragraph.strip()
    ]
    chunks: list[str] = []
    current = ""

    for paragraph in paragraphs:
        parts = _split_long_text(paragraph, safe_max)
        for part in parts:
            if not current:
                current = part
            elif len(current) + 2 + len(part) <= safe_max:
                current = f"{current}\n\n{part}"
            else:
                chunks.append(current)
                current = _overlap_prefix(current, safe_overlap, part)
    if current:
        chunks.append(current)
    return chunks


def chunk_segments(
    segments: list[ExtractedSegment],
    *,
    max_chars: int = DEFAULT_CHUNK_CHARS,
    overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> list[ExtractedSegment]:
    """Chunk each locator boundary independently while preserving provenance.

    Text from different pages, slides, or sheets is never merged into one chunk;
    doing so would make a citation point at a location containing only part of
    the retrieved passage. Long individual segments retain the usual overlap.
    """

    chunks: list[ExtractedSegment] = []
    for segment in segments:
        for text in chunk_text(segment.text, max_chars=max_chars, overlap=overlap):
            chunks.append(
                ExtractedSegment(
                    text=text,
                    page_start=segment.page_start,
                    page_end=segment.page_end,
                    locator=segment.locator,
                )
            )
    return chunks


def _extension(filename: str) -> str:
    if "." not in filename:
        return ""
    return f".{filename.rsplit('.', 1)[-1].lower()}"


def _decode_text(content: bytes) -> str:
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return content.decode("latin-1", errors="replace")


def _extract_docx_segments(source: BinarySource) -> list[ExtractedSegment]:
    try:
        _rewind_source(source)
        archive_source = BytesIO(source) if isinstance(source, bytes) else source
        with zipfile.ZipFile(archive_source) as archive:
            document_xml = archive.read("word/document.xml")
    except (KeyError, OSError, zipfile.BadZipFile):
        return []
    try:
        root = ElementTree.fromstring(document_xml)
    except (ElementTree.ParseError, DefusedXmlException):
        return []

    body = next((element for element in root.iter() if _local_name(element.tag) == "body"), None)
    if body is None:
        return []

    blocks: list[str] = []
    for element in body:
        name = _local_name(element.tag)
        if name == "p":
            text = _word_paragraph_text(element)
            if text:
                blocks.append(text)
        elif name == "tbl":
            table_text = _word_table_text(element)
            if table_text:
                blocks.append(table_text)
    if not blocks:
        return []
    return [ExtractedSegment(text="\n\n".join(blocks))]


def _word_paragraph_text(paragraph: ElementTree.Element) -> str:
    parts: list[str] = []
    for node in paragraph.iter():
        name = _local_name(node.tag)
        if name == "t":
            parts.append(node.text or "")
        elif name == "tab":
            parts.append("\t")
        elif name in {"br", "cr"}:
            parts.append("\n")
    return "".join(parts).strip()


def _word_table_text(table: ElementTree.Element) -> str:
    rows: list[str] = []
    for row in table:
        if _local_name(row.tag) != "tr":
            continue
        cells: list[str] = []
        for cell in row:
            if _local_name(cell.tag) != "tc":
                continue
            paragraphs = [
                _word_paragraph_text(paragraph)
                for paragraph in cell.iter()
                if _local_name(paragraph.tag) == "p"
            ]
            cells.append(" / ".join(text for text in paragraphs if text))
        while cells and not cells[-1]:
            cells.pop()
        if cells:
            rows.append(" | ".join(cells))
    return "\n\n".join(rows)


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _extract_xlsx_segments(source: BinarySource) -> list[ExtractedSegment]:
    try:
        from openpyxl import load_workbook  # type: ignore[import-not-found]
    except ImportError:
        return []

    workbook = None
    try:
        workbook = load_workbook(
            BytesIO(_read_all_source_bytes(source)),
            read_only=True,
            data_only=False,
            keep_links=False,
        )
        segments: list[ExtractedSegment] = []
        for worksheet in workbook.worksheets:
            rows: list[str] = []
            for values in worksheet.iter_rows(values_only=True):
                cells = [_spreadsheet_cell_text(value) for value in values]
                while cells and not cells[-1]:
                    cells.pop()
                if cells:
                    rows.append(" | ".join(cells))
            text = "\n\n".join(rows)
            if text:
                segments.append(ExtractedSegment(text=text, locator=f"Sheet: {worksheet.title}"))
        return segments
    except Exception:
        return []
    finally:
        if workbook is not None:
            workbook.close()


def _spreadsheet_cell_text(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    return str(value).strip()


def _extract_pptx_segments(source: BinarySource) -> list[ExtractedSegment]:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError:
        return []

    try:
        presentation = Presentation(BytesIO(_read_all_source_bytes(source)))
    except Exception:
        return []

    segments: list[ExtractedSegment] = []
    for index, slide in enumerate(presentation.slides, start=1):
        blocks: list[str] = []
        for shape in slide.shapes:
            blocks.extend(_pptx_shape_blocks(shape))
        text = "\n\n".join(block for block in blocks if block.strip())
        if text:
            segments.append(ExtractedSegment(text=text, locator=f"Slide {index}"))
    return segments


def _pptx_shape_blocks(shape: object) -> list[str]:
    blocks: list[str] = []
    if bool(getattr(shape, "has_table", False)):
        table = getattr(shape, "table", None)
        rows = getattr(table, "rows", []) if table is not None else []
        rendered_rows: list[str] = []
        for row in rows:
            cells = [str(getattr(cell, "text", "") or "").strip() for cell in row.cells]
            while cells and not cells[-1]:
                cells.pop()
            if cells:
                rendered_rows.append(" | ".join(cells))
        if rendered_rows:
            blocks.append("\n\n".join(rendered_rows))
    elif bool(getattr(shape, "has_text_frame", False)):
        text = str(getattr(shape, "text", "") or "").strip()
        if text:
            blocks.append(text)

    nested_shapes = getattr(shape, "shapes", None)
    if nested_shapes is not None:
        for nested in nested_shapes:
            blocks.extend(_pptx_shape_blocks(nested))
    return blocks


def _extract_eml_segments(source: BinarySource) -> list[ExtractedSegment]:
    try:
        message = BytesParser(policy=policy.default).parsebytes(_read_all_source_bytes(source))
    except Exception:
        return []

    headers = _email_header_lines(
        {
            "From": message.get("From"),
            "To": message.get("To"),
            "Cc": message.get("Cc"),
            "Date": message.get("Date"),
            "Subject": message.get("Subject"),
        }
    )
    body = _email_message_body(message)
    text = "\n\n".join(part for part in (headers, body) if part)
    return [ExtractedSegment(text=text)] if text else []


def _email_message_body(message: object) -> str:
    get_body = getattr(message, "get_body", None)
    body_part = get_body(preferencelist=("plain", "html")) if callable(get_body) else None
    if body_part is not None:
        try:
            body = body_part.get_content()
        except (KeyError, LookupError, UnicodeError):
            body = ""
        if not isinstance(body, str):
            return ""
        if body_part.get_content_type() == "text/html":
            return _extract_html_text(body)
        return body

    if not bool(getattr(message, "is_multipart", lambda: False)()):
        try:
            body = message.get_content()
        except (AttributeError, KeyError, LookupError, UnicodeError):
            return ""
        return body if isinstance(body, str) else ""
    return ""


def _extract_msg_segments(source: BinarySource) -> list[ExtractedSegment]:
    try:
        from oxmsg import Message  # type: ignore[import-not-found]
    except ImportError:
        return []

    content = _read_all_source_bytes(source)
    if not content:
        return []
    try:
        message = Message.load(content)
    except Exception:
        return []

    raw_headers = getattr(message, "message_headers", {})
    headers_by_name = (
        {str(name).casefold(): value for name, value in raw_headers.items()}
        if isinstance(raw_headers, dict)
        else {}
    )
    sent_date = getattr(message, "sent_date", None)
    headers = _email_header_lines(
        {
            "From": getattr(message, "sender", None) or headers_by_name.get("from"),
            "To": headers_by_name.get("to"),
            "Cc": headers_by_name.get("cc"),
            "Date": sent_date.isoformat() if sent_date is not None else headers_by_name.get("date"),
            "Subject": getattr(message, "subject", None) or headers_by_name.get("subject"),
        }
    )
    plain_body = getattr(message, "body", None)
    if plain_body is not None:
        body = str(plain_body)
    else:
        body = ""
        html_body = getattr(message, "html_body", None)
        if isinstance(html_body, bytes):
            html_body = _decode_text(html_body)
        if isinstance(html_body, str):
            body = _extract_html_text(html_body)
    text = "\n\n".join(part for part in (headers, body) if part)
    return [ExtractedSegment(text=text)] if text else []


def _email_header_lines(values: dict[str, object]) -> str:
    lines: list[str] = []
    for label, raw_value in values.items():
        if raw_value is None:
            continue
        if isinstance(raw_value, (list, tuple, set)):
            value = ", ".join(str(item).strip() for item in raw_value if str(item).strip())
        else:
            value = str(raw_value).strip()
        if value:
            lines.append(f"{label}: {value}")
    return "\n".join(lines)


def _extract_html_text(text: str) -> str:
    without_scripts = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", text)
    without_tags = re.sub(r"(?s)<[^>]+>", " ", without_scripts)
    return html.unescape(without_tags)


def _extract_rtf_text(text: str) -> str:
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
    text = text.replace("{", " ").replace("}", " ").replace("\\", " ")
    return text


def _extract_pdf_segments(
    source: BinarySource,
    *,
    ocr_enabled: bool,
    ocr_max_pages: int,
    ocr_page_timeout_seconds: float,
) -> list[ExtractedSegment]:
    try:
        from pypdf import PdfReader  # type: ignore[import-not-found]
    except ImportError:
        return []
    try:
        _rewind_source(source)
        reader_source = BytesIO(source) if isinstance(source, bytes) else source
        reader = PdfReader(reader_source)
        page_texts = [_normalize_text(page.extract_text() or "") for page in reader.pages]
    except Exception:
        return []

    if ocr_enabled:
        pages_needing_ocr = [
            index
            for index, text in enumerate(page_texts[:ocr_max_pages])
            if len(text) < MIN_PDF_PAGE_TEXT_CHARS
        ]
        if pages_needing_ocr:
            ocr_text = _ocr_pdf_pages(
                source,
                pages_needing_ocr,
                timeout_seconds=ocr_page_timeout_seconds,
            )
            for page_index, text in ocr_text.items():
                if text:
                    page_texts[page_index] = text

    return [
        ExtractedSegment(
            text=text,
            page_start=index,
            page_end=index,
            locator=f"Page {index}",
        )
        for index, text in enumerate(page_texts, start=1)
        if text
    ]


def _ocr_pdf_pages(
    source: BinarySource,
    page_indexes: list[int],
    *,
    timeout_seconds: float,
) -> dict[int, str]:
    try:
        import pypdfium2 as pdfium  # type: ignore[import-not-found]
    except ImportError:
        return {}

    _rewind_source(source)
    pdf_source = source if isinstance(source, bytes) else source
    try:
        pdf = pdfium.PdfDocument(pdf_source)
    except Exception:
        return {}

    results: dict[int, str] = {}
    try:
        for page_index in page_indexes:
            if page_index >= len(pdf):
                continue
            page = pdf[page_index]
            bitmap = None
            image = None
            try:
                bitmap = page.render(scale=2, rotation=0)
                image = bitmap.to_pil()
                results[page_index] = _ocr_image(
                    image,
                    timeout_seconds=timeout_seconds,
                )
            except Exception:
                continue
            finally:
                if image is not None:
                    image.close()
                if bitmap is not None:
                    bitmap.close()
                page.close()
    finally:
        pdf.close()
    return results


def _extract_image_segments(
    source: BinarySource,
    *,
    ocr_max_pages: int,
    ocr_page_timeout_seconds: float,
) -> list[ExtractedSegment]:
    try:
        from PIL import Image, ImageSequence  # type: ignore[import-not-found]
    except ImportError:
        return []

    _rewind_source(source)
    image_source = BytesIO(source) if isinstance(source, bytes) else source
    try:
        with Image.open(image_source) as image:
            segments: list[ExtractedSegment] = []
            for index, frame in enumerate(ImageSequence.Iterator(image)):
                if index >= ocr_max_pages:
                    break
                frame_copy = frame.copy()
                try:
                    text = _ocr_image(
                        frame_copy,
                        timeout_seconds=ocr_page_timeout_seconds,
                    )
                    if text:
                        page_number = index + 1
                        segments.append(
                            ExtractedSegment(
                                text=text,
                                page_start=page_number,
                                page_end=page_number,
                                locator=f"Page {page_number}",
                            )
                        )
                finally:
                    frame_copy.close()
    except Exception:
        return []
    return segments


def _ocr_image(image: object, *, timeout_seconds: float) -> str:
    try:
        import pytesseract  # type: ignore[import-not-found]
        from PIL import ImageOps  # type: ignore[import-not-found]
    except ImportError:
        return ""

    try:
        prepared = ImageOps.exif_transpose(image).convert("L")
        prepared = ImageOps.autocontrast(prepared)
        try:
            return _normalize_text(
                pytesseract.image_to_string(
                    prepared,
                    lang="eng",
                    config="--oem 1 --psm 3",
                    timeout=timeout_seconds,
                )
            )
        finally:
            prepared.close()
    except Exception:
        return ""


def _read_source_bytes(source: BinarySource, *, max_chars: int) -> bytes:
    if isinstance(source, bytes):
        return source[: max_chars * 4 + 4]
    _rewind_source(source)
    return source.read(max_chars * 4 + 4)


def _read_all_source_bytes(source: BinarySource) -> bytes:
    if isinstance(source, bytes):
        return source
    _rewind_source(source)
    return source.read()


def _rewind_source(source: BinarySource) -> None:
    if not isinstance(source, bytes):
        source.seek(0)


def _limit_segments(
    segments: list[ExtractedSegment],
    *,
    max_chars: int = MAX_EXTRACTED_TEXT_CHARS,
) -> list[ExtractedSegment]:
    """Normalize segments under one global joined-text character budget."""

    remaining = max(0, max_chars)
    limited: list[ExtractedSegment] = []
    for segment in segments:
        normalized = _normalize_text(segment.text)
        if not normalized:
            continue
        separator_chars = 2 if limited else 0
        if remaining <= separator_chars:
            break
        remaining -= separator_chars
        text = normalized[:remaining].rstrip()
        if not text:
            break
        limited.append(
            ExtractedSegment(
                text=text,
                page_start=segment.page_start,
                page_end=segment.page_end,
                locator=segment.locator,
            )
        )
        remaining -= len(text)
        if remaining <= 0:
            break
    return limited


def _segments_text(segments: list[ExtractedSegment]) -> str | None:
    text = "\n\n".join(segment.text for segment in segments if segment.text)
    return text or None


def _normalize_text(text: str) -> str:
    lines = [" ".join(line.replace("\x00", " ").split()) for line in text.splitlines()]
    paragraphs: list[str] = []
    current: list[str] = []
    for line in lines:
        if line:
            current.append(line)
        elif current:
            paragraphs.append(" ".join(current))
            current = []
    if current:
        paragraphs.append(" ".join(current))
    return "\n\n".join(paragraphs).strip()


def _split_long_text(text: str, max_chars: int) -> list[str]:
    if len(text) <= max_chars:
        return [text]
    sentences = re.split(r"(?<=[.!?])\s+", text)
    parts: list[str] = []
    current = ""
    for sentence in sentences:
        if not sentence:
            continue
        for segment in _split_oversized_sentence(sentence, max_chars):
            if not current:
                current = segment
            elif len(current) + 1 + len(segment) <= max_chars:
                current = f"{current} {segment}"
            else:
                parts.append(current)
                current = segment
    if current:
        parts.append(current)
    return parts


def _split_oversized_sentence(sentence: str, max_chars: int) -> list[str]:
    if len(sentence) <= max_chars:
        return [sentence]

    segments: list[str] = []
    current = ""
    for word in sentence.split():
        if len(word) > max_chars:
            if current:
                segments.append(current)
                current = ""
            segments.extend(
                word[offset : offset + max_chars] for offset in range(0, len(word), max_chars)
            )
        elif not current:
            current = word
        elif len(current) + 1 + len(word) <= max_chars:
            current = f"{current} {word}"
        else:
            segments.append(current)
            current = word
    if current:
        segments.append(current)
    return segments


def _overlap_prefix(previous: str, overlap: int, next_text: str) -> str:
    if overlap <= 0:
        return next_text
    prefix = previous[-overlap:].strip()
    if not prefix:
        return next_text
    return f"{prefix}\n\n{next_text}"
